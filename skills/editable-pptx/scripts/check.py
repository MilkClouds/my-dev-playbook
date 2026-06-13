#!/usr/bin/env python3
"""Structural QA for a .pptx (no rendering): per-slide shape/text inventory plus a
cheap text-box overlap detector to catch layout bugs before opening PowerPoint.

Usage:  python check.py deck.pptx
Requires (in your QA venv):  python-pptx
"""
import sys
from pptx import Presentation
from pptx.util import Emu

EMU = 914400  # per inch


def overlaps(a, b):
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    ix = max(0, min(ax + aw, bx + bw) - max(ax, bx))
    iy = max(0, min(ay + ah, by + bh) - max(ay, by))
    inter = ix * iy
    if inter <= 0:
        return 0.0
    return inter / min(aw * ah, bw * bh)


def main(path):
    p = Presentation(path)
    sw, sh = Emu(p.slide_width).inches, Emu(p.slide_height).inches
    print(f"{path}: {len(p.slides)} slides @ {sw:.2f} x {sh:.2f} in")
    flagged = 0
    for i, sl in enumerate(p.slides, 1):
        boxes = []  # only text boxes with content (decorative rects/lines ignored)
        nshapes = nchart = ntable = 0
        for shp in sl.shapes:
            nshapes += 1
            if shp.has_chart:
                nchart += 1
            if shp.has_table:
                ntable += 1
            if shp.has_text_frame and shp.text_frame.text.strip():
                try:
                    boxes.append((Emu(shp.left).inches, Emu(shp.top).inches,
                                  Emu(shp.width).inches, Emu(shp.height).inches,
                                  shp.text_frame.text.strip()[:28]))
                except (TypeError, ValueError):
                    pass
        # off-canvas / overflow check
        for (x, y, w, h, t) in boxes:
            if x < -0.05 or y < -0.05 or x + w > sw + 0.05 or y + h > sh + 0.05:
                print(f"  [slide {i}] OFF-CANVAS  '{t}'  ({x:.1f},{y:.1f},{w:.1f}x{h:.1f})")
                flagged += 1
        # text-box overlap check (>35% of the smaller box)
        for m in range(len(boxes)):
            for n in range(m + 1, len(boxes)):
                ov = overlaps(boxes[m][:4], boxes[n][:4])
                if ov > 0.35:
                    print(f"  [slide {i}] OVERLAP {ov:.0%}  '{boxes[m][4]}'  <->  '{boxes[n][4]}'")
                    flagged += 1
        print(f"  slide {i:2d}: {nshapes:2d} shapes · {len(boxes)} text · chart={nchart} table={ntable}")
    print(f"done — {flagged} layout warning(s)")


if __name__ == "__main__":
    if len(sys.argv) != 2:
        sys.exit("usage: check.py deck.pptx")
    main(sys.argv[1])
